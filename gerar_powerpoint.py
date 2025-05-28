import pandas as pd
from pptx import Presentation

def gerar():
    df = pd.read_excel("indicadores.xlsx")
    dados = dict(zip(df["Indicador"], df["Valor"]))

    for chave, valor in dados.items():
        chave_lower = chave.lower()
        try:
            valor_numerico = float(str(valor).replace(".", "").replace(",", "."))
        except ValueError:
            continue

        if "r$" in chave_lower:
            dados[chave] = f"R$ {valor_numerico:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
        elif "num" in chave_lower:
            dados[chave] = f"{int(valor_numerico):,}".replace(",", ".")

    def substituir_texto(ppt, dados):
        for slide in ppt.slides:
            for forma in slide.shapes:
                if forma.has_text_frame:
                    for paragrafo in forma.text_frame.paragraphs:
                        for run in paragrafo.runs:
                            for chave, valor in dados.items():
                                marcador = f"{{{{{chave}}}}}"
                                if marcador in run.text:
                                    run.text = run.text.replace(marcador, str(valor))

    # Lista dos arquivos modelo
    modelos = [
        ("slide_padrão_medicina.pptx", "preenchido_medicina.pptx"),
        ("slide_padrão_odontologia.pptx", "preenchido_odontologia.pptx"),
    ]

    for modelo, destino in modelos:
        ppt = Presentation(modelo)
        substituir_texto(ppt, dados)
        ppt.save(destino)

gerar()